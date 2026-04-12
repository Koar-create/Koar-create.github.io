import pptx
import os

path = r"G:\Other computers\My Laptop\slides\slide.Apr10th2026 ( ! ).pptx"

if not os.path.exists(path):
    print(f"File not found: {path}")
else:
    prj = pptx.Presentation(path)
    for i, slide in enumerate(prj.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
