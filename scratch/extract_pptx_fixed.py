import pptx
import os
import sys

# Set output encoding to UTF-8
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

path = r"G:\Other computers\My Laptop\slides\slide.Apr10th2026 ( ! ).pptx"

if not os.path.exists(path):
    print(f"File not found: {path}")
else:
    prj = pptx.Presentation(path)
    for i, slide in enumerate(prj.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
